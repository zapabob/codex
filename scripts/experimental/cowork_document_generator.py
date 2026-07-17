#!/usr/bin/env python3
"""
ClaudeCowork-style Document Generation Engine
Excel（数式対応）、Word（書式設定）、PowerPoint生成機能
"""

import logging
import os
from pathlib import Path
from typing import Dict, List, Optional, Any
from datetime import datetime
import json

try:
    from openpyxl import Workbook
    from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
    from openpyxl.chart import BarChart, LineChart, PieChart
    from openpyxl.utils import get_column_letter
    from docx import Document
    from docx.shared import Pt, RGBColor, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError as e:
    print(f"必要なライブラリがインストールされていません: {e}")
    print("以下のコマンドでインストールしてください:")
    print("pip install openpyxl python-docx python-pptx")
    exit(1)

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class DocumentGenerationEngine:
    """
    ClaudeCoworkスタイルのドキュメント生成エンジン

    機能:
    - Excel生成（数式、グラフ、書式設定）
    - Word生成（スタイル、目次、画像挿入）
    - PowerPoint生成（テンプレート、アニメーション）
    """

    def __init__(self):
        self.logger = logging.getLogger("DocumentGenerationEngine")

    def generate_excel(self, output_path: str, data: Dict[str, Any]) -> Dict[str, Any]:
        """
        Excelスプレッドシート生成

        Args:
            output_path: 出力ファイルパス
            data: 生成データ
                - sheets: シート定義のリスト
                - formulas: 数式定義
                - charts: グラフ定義
                - styles: スタイル定義
        """
        try:
            wb = Workbook()
            wb.remove(wb.active)  # デフォルトシート削除

            for sheet_data in data.get("sheets", []):
                sheet_name = sheet_data.get("name", "Sheet1")
                ws = wb.create_sheet(title=sheet_name)

                # データ挿入
                rows = sheet_data.get("rows", [])
                for row_idx, row_data in enumerate(rows, start=1):
                    for col_idx, cell_value in enumerate(row_data, start=1):
                        cell = ws.cell(row=row_idx, column=col_idx)
                        cell.value = cell_value

                # 数式適用
                formulas = sheet_data.get("formulas", {})
                for cell_ref, formula in formulas.items():
                    ws[cell_ref].value = formula

                # スタイル適用
                styles = sheet_data.get("styles", {})
                self._apply_excel_styles(ws, styles)

                # グラフ生成
                charts = sheet_data.get("charts", [])
                for chart_data in charts:
                    self._create_excel_chart(ws, chart_data)

                # 列幅自動調整
                for column in ws.columns:
                    max_length = 0
                    column_letter = get_column_letter(column[0].column)
                    for cell in column:
                        try:
                            if len(str(cell.value)) > max_length:
                                max_length = len(str(cell.value))
                        except:
                            pass
                    adjusted_width = min(max_length + 2, 50)
                    ws.column_dimensions[column_letter].width = adjusted_width

            wb.save(output_path)
            self.logger.info(f"Excel生成完了: {output_path}")
            return {"success": True, "path": output_path}

        except Exception as e:
            self.logger.error(f"Excel生成エラー: {e}")
            return {"success": False, "error": str(e)}

    def _apply_excel_styles(self, ws, styles: Dict[str, Any]):
        """Excelスタイル適用"""
        # ヘッダースタイル
        if "header" in styles:
            header_style = styles["header"]
            header_row = header_style.get("row", 1)

            fill = PatternFill(
                start_color=header_style.get("bg_color", "366092"),
                end_color=header_style.get("bg_color", "366092"),
                fill_type="solid",
            )
            font = Font(bold=True, color=header_style.get("text_color", "FFFFFF"))

            for cell in ws[header_row]:
                cell.fill = fill
                cell.font = font
                cell.alignment = Alignment(horizontal="center", vertical="center")

        # ボーダー
        if "border" in styles:
            thin_border = Border(
                left=Side(style="thin"),
                right=Side(style="thin"),
                top=Side(style="thin"),
                bottom=Side(style="thin"),
            )
            for row in ws.iter_rows():
                for cell in row:
                    cell.border = thin_border

    def _create_excel_chart(self, ws, chart_data: Dict[str, Any]):
        """Excelグラフ生成"""
        chart_type = chart_data.get("type", "bar")
        data_range = chart_data.get("data_range")
        title = chart_data.get("title", "Chart")

        if chart_type == "bar":
            chart = BarChart()
        elif chart_type == "line":
            chart = LineChart()
        elif chart_type == "pie":
            chart = PieChart()
        else:
            chart = BarChart()

        chart.title = title
        chart.style = chart_data.get("style", 10)

        # データ範囲設定
        if data_range:
            data = ws[data_range]
            chart.add_data(data, titles_from_data=True)
            ws.add_chart(chart, chart_data.get("position", "E2"))

    def generate_word(
        self, output_path: str, content: Dict[str, Any]
    ) -> Dict[str, Any]:
        """
        Word文書生成

        Args:
            output_path: 出力ファイルパス
            content: 生成コンテンツ
                - title: タイトル
                - sections: セクションリスト
                - styles: スタイル定義
                - images: 画像パスリスト
        """
        try:
            doc = Document()

            # タイトル
            if "title" in content:
                title = doc.add_heading(content["title"], 0)
                title.alignment = WD_ALIGN_PARAGRAPH.CENTER

            # セクション追加
            for section in content.get("sections", []):
                # 見出し
                if "heading" in section:
                    doc.add_heading(section["heading"], level=section.get("level", 1))

                # 段落
                if "paragraphs" in section:
                    for para_text in section["paragraphs"]:
                        para = doc.add_paragraph(para_text)

                        # スタイル適用
                        if "style" in section:
                            style = section["style"]
                            if "bold" in style and style["bold"]:
                                for run in para.runs:
                                    run.bold = True
                            if "italic" in style and style["italic"]:
                                for run in para.runs:
                                    run.italic = True
                            if "color" in style:
                                color = self._parse_color(style["color"])
                                for run in para.runs:
                                    run.font.color.rgb = color

                # リスト
                if "list" in section:
                    list_type = section["list"].get("type", "bullet")
                    items = section["list"].get("items", [])

                    for item in items:
                        if list_type == "bullet":
                            doc.add_paragraph(item, style="List Bullet")
                        else:
                            doc.add_paragraph(item, style="List Number")

                # 画像
                if "images" in section:
                    for image_path in section["images"]:
                        if Path(image_path).exists():
                            doc.add_picture(image_path, width=Inches(5))

            # 目次（オプション）
            if content.get("table_of_contents", False):
                # 目次は手動で追加する必要がある（python-docxの制限）
                doc.add_paragraph("目次（手動で更新してください）")

            doc.save(output_path)
            self.logger.info(f"Word生成完了: {output_path}")
            return {"success": True, "path": output_path}

        except Exception as e:
            self.logger.error(f"Word生成エラー: {e}")
            return {"success": False, "error": str(e)}

    def _parse_color(self, color_str: str) -> RGBColor:
        """色文字列をRGBColorに変換"""
        if color_str.startswith("#"):
            color_str = color_str[1:]
        r = int(color_str[0:2], 16)
        g = int(color_str[2:4], 16)
        b = int(color_str[4:6], 16)
        return RGBColor(r, g, b)

    def generate_powerpoint(
        self, output_path: str, presentation_data: Dict[str, Any]
    ) -> Dict[str, Any]:
        """
        PowerPointプレゼンテーション生成

        Args:
            output_path: 出力ファイルパス
            presentation_data: プレゼンテーションデータ
                - title: タイトル
                - slides: スライドリスト
                - template: テンプレートパス（オプション）
        """
        try:
            # テンプレート読み込み（オプション）
            if (
                "template" in presentation_data
                and Path(presentation_data["template"]).exists()
            ):
                prs = Presentation(presentation_data["template"])
            else:
                prs = Presentation()

            # タイトルスライド
            if "title" in presentation_data:
                title_slide_layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(title_slide_layout)
                title = slide.shapes.title
                subtitle = slide.placeholders[1]

                title.text = presentation_data["title"]
                if "subtitle" in presentation_data:
                    subtitle.text = presentation_data["subtitle"]

            # スライド追加
            for slide_data in presentation_data.get("slides", []):
                slide_layout = prs.slide_layouts[slide_data.get("layout", 1)]
                slide = prs.slides.add_slide(slide_layout)

                # タイトル
                if slide.shapes.title:
                    slide.shapes.title.text = slide_data.get("title", "")

                # コンテンツ
                if "content" in slide_data:
                    content_placeholder = (
                        slide.placeholders[1] if len(slide.placeholders) > 1 else None
                    )
                    if content_placeholder:
                        tf = content_placeholder.text_frame
                        tf.text = slide_data["content"]

                        # 段落追加
                        if "paragraphs" in slide_data:
                            for para_text in slide_data["paragraphs"]:
                                p = tf.add_paragraph()
                                p.text = para_text
                                p.level = slide_data.get("paragraph_level", 0)

                # 画像
                if "images" in slide_data:
                    for img_data in slide_data["images"]:
                        img_path = img_data.get("path")
                        if img_path and Path(img_path).exists():
                            left = Inches(img_data.get("left", 1))
                            top = Inches(img_data.get("top", 2))
                            width = Inches(img_data.get("width", 5))
                            height = Inches(img_data.get("height", 3))
                            slide.shapes.add_picture(img_path, left, top, width, height)

            prs.save(output_path)
            self.logger.info(f"PowerPoint生成完了: {output_path}")
            return {"success": True, "path": output_path}

        except Exception as e:
            self.logger.error(f"PowerPoint生成エラー: {e}")
            return {"success": False, "error": str(e)}


def main():
    """テスト実行"""
    engine = DocumentGenerationEngine()

    # Excel生成テスト
    excel_data = {
        "sheets": [
            {
                "name": "Sales",
                "rows": [
                    ["Product", "Q1", "Q2", "Q3", "Q4"],
                    ["Product A", 100, 120, 130, 140],
                    ["Product B", 200, 210, 220, 230],
                    ["Product C", 150, 160, 170, 180],
                ],
                "formulas": {
                    "B5": "=SUM(B2:B4)",
                    "C5": "=SUM(C2:C4)",
                    "D5": "=SUM(D2:D4)",
                    "E5": "=SUM(E2:E4)",
                },
                "styles": {
                    "header": {"row": 1, "bg_color": "366092", "text_color": "FFFFFF"}
                },
                "charts": [
                    {
                        "type": "bar",
                        "data_range": "A1:E4",
                        "title": "Sales by Quarter",
                        "position": "A7",
                    }
                ],
            }
        ]
    }

    result = engine.generate_excel("test_output.xlsx", excel_data)
    print(f"Excel生成結果: {result}")

    # Word生成テスト
    word_content = {
        "title": "サンプル文書",
        "sections": [
            {
                "heading": "はじめに",
                "level": 1,
                "paragraphs": [
                    "これはサンプル文書です。",
                    "ClaudeCoworkスタイルのドキュメント生成機能をテストしています。",
                ],
                "style": {"bold": False},
            },
            {
                "heading": "機能一覧",
                "level": 1,
                "list": {
                    "type": "bullet",
                    "items": ["Excel生成", "Word生成", "PowerPoint生成"],
                },
            },
        ],
        "table_of_contents": False,
    }

    result = engine.generate_word("test_output.docx", word_content)
    print(f"Word生成結果: {result}")

    # PowerPoint生成テスト
    ppt_data = {
        "title": "サンプルプレゼンテーション",
        "subtitle": "ClaudeCowork統合テスト",
        "slides": [
            {
                "title": "スライド1",
                "content": "最初のスライド",
                "paragraphs": ["ポイント1", "ポイント2", "ポイント3"],
            },
            {
                "title": "スライド2",
                "content": "2番目のスライド",
                "paragraphs": ["詳細情報1", "詳細情報2"],
            },
        ],
    }

    result = engine.generate_powerpoint("test_output.pptx", ppt_data)
    print(f"PowerPoint生成結果: {result}")


if __name__ == "__main__":
    main()
